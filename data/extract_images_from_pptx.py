from pptx import Presentation
import os
import subprocess

def extract_images_as_png(pptx_file, output_dir):
    prs = Presentation(pptx_file)
    os.makedirs(output_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                ext = image.ext  # e.g. 'jpeg', 'png', 'wmf', 'emf'
                raw_path = os.path.join(output_dir, f"{i}.{ext}")
                png_path = os.path.join(output_dir, f"{i}.png")

                # save original blob
                with open(raw_path, "wb") as f:
                    f.write(image.blob)

                if ext.lower() in ["wmf", "emf"]:
                    # convert with ImageMagick
                    subprocess.run(
                        ["magick", raw_path, png_path],
                        check=True
                    )
                    os.remove(raw_path)  # cleanup original WMF/EMF
                else:
                    # already a bitmap, just rename/copy to png if needed
                    if ext.lower() != "png":
                        subprocess.run(
                            ["magick", raw_path, png_path],
                            check=True
                        )
                        os.remove(raw_path)
                break  # only first image per slide

# Example usage
extract_images_as_png("in/DO-80 stimuli.pptx.odp", "out/do80")
